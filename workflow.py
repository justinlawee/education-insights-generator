from typing import Dict, TypedDict, Optional, List, Tuple, Any
from collections import Counter

from langgraph.graph import StateGraph, START, END
from langchain_google_vertexai import ChatVertexAI
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_google_community import VertexAISearchRetriever

import prompts
import settings
from google.cloud import bigquery
import bq_functions
import utils
import json
import os
import io
import plotly.io  # For converting Plotly fig to JSON
from docx import Document  # For .docx files
import fitz  # PyMuPDF for .pdf files
from pptx import Presentation  # For .pptx files
import numpy as np  # For correlation functions

import pandas as pd
from langchain_core.messages import AIMessage, SystemMessage, HumanMessage


# TODO: Seeing about having a new agent to clarify questions and plan the workflow
# TODO: Gemini suggested adding an ethical and bias reviewer
# TODO: Updating the prompts
# TODO: Seeing if there is a use of a system prompt
# TODO: Seeing if there is a use to adding some memory to various agents
# TODO: Adding the insights to use the data as well as the graph
# TODO: Test out each of the various agents and temperatures

class AgentState(TypedDict):
    question: str
    database_schemas: str

    # General config (used by validator loops)
    max_num_retries_debug: int

    # LLM Configurations
    llm_config_sql_writer: Dict[str, Any]
    llm_config_sql_validator: Dict[str, Any]
    llm_config_sql_judge: Dict[str, Any]
    llm_config_bi_expert: Dict[str, Any]
    llm_config_bi_judge: Dict[str, Any]
    llm_config_python_writer: Dict[str, Any]
    llm_config_python_validator: Dict[str, Any]
    llm_config_python_judge: Dict[str, Any]
    llm_config_insights_generator: Dict[str, Any]
    llm_config_insights_judge: Dict[str, Any]

    # SQL Writer Agent State
    sql_writer_n_runs: int
    sql_writer_candidate_queries: List[str]

    sql_writer_generation_input_tokens: int
    sql_writer_generation_output_tokens: int

    # SQL Validator Agent State
    # Stores {'original_query': str, 'final_query': str, 'passed': bool, 'error': Optional[str], 'attempts': int}
    sql_validation_results: List[Dict[str, Any]]
    sql_passed_validation_queries: List[str]  # List of query *strings* that passed
    sql_validator_input_tokens: int
    sql_validator_output_tokens: int

    # SQL Judge Agent State
    query: str  # This will store the final query *selected* by the judge (or the single validated one)
    sql_writer_judge_reasoning: Optional[str]
    sql_writer_judge_input_tokens: int
    sql_writer_judge_output_tokens: int

    # SQL Executor Agent State
    df: pd.DataFrame  # Populated by execute_selected_sql_node
    sql_execution_error: Optional[str]  # Error during final execution

    # BI Expert Agent State
    visualization_request: str  # Selected visualization request
    bi_expert_n_runs: int
    bi_expert_candidate_requests: List[str]
    bi_expert_judge_reasoning: Optional[str]

    bi_expert_generation_input_tokens: int
    bi_expert_generation_output_tokens: int
    bi_expert_judge_input_tokens: int
    bi_expert_judge_output_tokens: int

    # Python Code Data Visualization Generator Agent State
    python_code_data_visualization: str  # Selected python code
    py_gen_n_runs: int
    py_gen_candidate_codes: List[str]
    py_gen_judge_reasoning: Optional[str]

    py_gen_generation_input_tokens: int
    py_gen_generation_output_tokens: int
    py_gen_judge_input_tokens: int
    py_gen_judge_output_tokens: int

    # Python Code Validator Agent State
    python_code_store_variables_dict: dict
    num_retries_debug_python_code_data_visualization: int
    result_debug_python_code_data_visualization: str
    error_msg_debug_python_code_data_visualization: str
    py_fixer_input_tokens: int
    py_fixer_output_tokens: int

    # Insights Agent State
    insights_n_runs: int
    insights_candidate_texts: List[str]
    insights_context_user_stories: str
    insights_context_mtss_docs: str
    insights_context_correlations: str
    generated_insights: str
    insights_generator_input_tokens: int
    insights_generator_output_tokens: int
    
    insights_judge_reasoning: Optional[str]
    insights_judge_input_tokens: int
    insights_judge_output_tokens: int


llm_sql_writer = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.5,
        project=settings.project_id,
        location="us-central1"
)

llm_sql_validator = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.2,
        project=settings.project_id,
        location="us-central1"
)

llm_sql_judge = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.1,
        project=settings.project_id,
        location="us-central1"
)

llm_bi_expert = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.5,
        project=settings.project_id,
        location="us-central1"
)

llm_bi_judge = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.1,
        project=settings.project_id,
        location="us-central1"
)

llm_python_writer = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.3,
        project=settings.project_id,
        location="us-central1"
)

llm_python_validator = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.2,
        project=settings.project_id,
        location="us-central1"
)

llm_python_judge = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.1,
        project=settings.project_id,
        location="us-central1"
)

llm_insights_generator = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",  # Or your preferred model
        temperature=0.4,  # Adjust as needed
        project=settings.project_id,
        location="us-central1"
)

llm_insights_judge = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",  # Or your preferred model
        temperature=0.1,  # Adjust as needed
        project=settings.project_id,
        location="us-central1"
)

llm_general = ChatVertexAI(
        model_name="gemini-2.0-flash-lite-001",
        temperature=0.3,
        project=settings.project_id,
        location="us-central1"
)

max_characters_error_msg_debug = 300

retriever = VertexAISearchRetriever(
        project_id=settings.project_id,
        location_id=settings.vertex_agent_builder_data_store_location,
        data_store_id=settings.vertex_agent_builder_data_store_id,
        max_documents=2,
        engine_data_type=1
)


# --------------------------------------------------------------------------
# --- CONTEXT LOADING FUNCTIONS (From User Snippets) ---
# --------------------------------------------------------------------------

def read_docx(file_path: str) -> str:
    """Reads text from a .docx file."""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
        return f"Error reading DOCX: {e}"


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts text from a .pdf file."""
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text("text") + "\n"
    except Exception as e:
        text = f"Error reading PDF {pdf_path}: {e}"
        print(text)
    return text


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extracts text from a .pptx file."""
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        text = f"Error reading PPTX {pptx_path}: {e}"
        print(text)
    return text


def extract_text_from_excel(excel_path: str) -> str:
    """Extracts text from an .xlsx or .xls file."""
    text = ""
    try:
        df_dict = pd.read_excel(excel_path, sheet_name=None)  # Read all sheets
        for sheet_name, sheet_data in df_dict.items():
            text += f"\n### Sheet: {sheet_name} ###\n"
            text += sheet_data.to_string(index=False) + "\n"
    except Exception as e:
        text = f"Error reading Excel {excel_path}: {e}"
        print(text)
    return text


def load_mtss_documentation_context(folder_path: str) -> str:
    """Loads context from various files in a folder."""
    file_contents = []
    if not os.path.isdir(folder_path):
        print(f"Error: MTSS Documentation folder not found: {folder_path}")
        return "Error: MTSS Documentation folder not found."

    for file in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file)
        if file.endswith(".pdf"):
            file_contents.append(f"### Document: {file} ###\n" + extract_text_from_pdf(file_path))
        elif file.endswith(".pptx"):
            file_contents.append(f"### Document: {file} ###\n" + extract_text_from_pptx(file_path))
        elif file.endswith(".xlsx") or file.endswith(".xls"):
            file_contents.append(f"### Document: {file} ###\n" + extract_text_from_excel(file_path))
        elif file.endswith(".docx"):
            file_contents.append(f"### Document: {file} ###\n" + read_docx(file_path))
    return "\n\n".join(file_contents) if file_contents else "No MTSS documents processed."


# --- Correlation Analysis Functions (From User Snippets, adapted slightly) ---
def load_and_merge_datasets(file1_path: str, file2_path: str) -> Optional[pd.DataFrame]:
    """Loads two datasets and merges them on 'Student_ID' and 'Year'."""
    try:
        df1 = pd.read_csv(file1_path) if file1_path.endswith('.csv') else pd.read_excel(file1_path)
        df2 = pd.read_csv(file2_path) if file2_path.endswith('.csv') else pd.read_excel(file2_path)
        merged_df = pd.merge(df1, df2, on=['Student_ID', 'Year'], how='inner')
        print("Correlation datasets loaded and merged successfully.")
        return merged_df
    except Exception as e:
        print(f"Error loading/merging correlation datasets: {e}")
        return None


def basic_info(df: pd.DataFrame):
    """Display basic information about the dataset."""
    print("Dataset Overview (Correlation Data):")
    # display(df.head()) # 'display' is for notebooks, use print for scripts
    print(df.head())
    print("\nShape of dataset:", df.shape)


def handle_missing_values(df: pd.DataFrame, strategy='mean') -> pd.DataFrame:
    """Handle missing values in numerical columns."""
    df_cleaned = df.copy()
    for col in df_cleaned.select_dtypes(include=[np.number]).columns:
        if df_cleaned[col].isnull().sum() > 0:
            if strategy == 'mean':
                df_cleaned[col].fillna(df_cleaned[col].mean(), inplace=True)
            elif strategy == 'median':
                df_cleaned[col].fillna(df_cleaned[col].median(), inplace=True)
            # Add other strategies if needed
    return df_cleaned


def encode_categorical(df: pd.DataFrame, max_categories=10) -> pd.DataFrame:
    """Encode categorical variables with low cardinality."""
    df_encoded = df.copy()
    id_columns = ['Student_Name', 'Student_ID', 'Current_School', 'Tested_School', 'Current_Homeroom']
    categorical_cols = [col for col in df_encoded.select_dtypes(include=['object']).columns if col not in id_columns]
    for col in categorical_cols:
        if df_encoded[col].nunique() <= max_categories:
            df_encoded = pd.get_dummies(df_encoded, columns=[col], drop_first=True)
        else:
            df_encoded.drop(columns=[col], inplace=True)  # Drop high-cardinality
    return df_encoded


def rename_merged_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Rename columns after merging, e.g., _x to _English, _y to _Math."""
    return df.rename(columns=lambda x: x.replace('_x', '_English').replace('_y', '_Math'))


def generate_correlation_matrix_string(df: pd.DataFrame, threshold=0.5) -> str:
    """Generate a string representation of strong correlations."""
    if df is None or df.empty:
        return "Correlation analysis could not be performed (No data)."
    try:
        corr_matrix = df.corr(numeric_only=True).dropna(how='all').dropna(axis=1, how='all')
        strong_corrs = corr_matrix[abs(corr_matrix) >= threshold]
        strong_corrs = strong_corrs.dropna(how='all').dropna(axis=1, how='all')

        if not strong_corrs.empty:
            # Optional: Display plot if in a notebook environment and libraries are imported
            # plt.figure(figsize=(12, 10))
            # sns.heatmap(strong_corrs, annot=True, cmap='coolwarm', fmt='.2f')
            # plt.title("Strong Correlations")
            # plt.show()
            return strong_corrs.to_string()
        else:
            return "No strong correlations found above the threshold."
    except Exception as e:
        print(f"Error generating correlation matrix string: {e}")
        return f"Error in correlation analysis: {e}"


def full_correlation_analysis_string(file1_path: str, file2_path: str) -> str:
    """Runs the full correlation analysis pipeline and returns a string summary."""
    df = load_and_merge_datasets(file1_path, file2_path)
    if df is not None:
        df = rename_merged_columns(df)
        # basic_info(df) # Optional: print basic info
        df = handle_missing_values(df, strategy='mean')
        df = encode_categorical(df)
        return generate_correlation_matrix_string(df)
    else:
        return "Correlation analysis could not be performed due to data loading issues."


# --- Helper function to parse judge's JSON output ---

def parse_judge_response(judge_response_content: str, num_unique_candidates: int) -> Tuple[
    Optional[int], Optional[str]]:
    try:
        # Attempt to find JSON block if the LLM wraps it with markdown
        match_json = utils.extract_code_block(judge_response_content, "json")
        if match_json:
            response_json = json.loads(match_json)
        else:  # Assume raw content is JSON or attempt direct parse
            response_json = json.loads(judge_response_content)

        selected_index = response_json.get("best_index")
        if selected_index is None:
            selected_index = response_json.get("best_query_index",
                                               response_json.get("selected_index"))  # Allow variations

        reasoning = response_json.get("reasoning", "")

        if isinstance(selected_index, str) and selected_index.isdigit():
            selected_index = int(selected_index)
        elif not isinstance(selected_index, int):
            print(f"Warning: Judge returned non-integer index: {selected_index}")
            return None, reasoning + f" (Error: Invalid index type '{type(selected_index)}' for index {selected_index})."

        if not (1 <= selected_index <= num_unique_candidates):
            print(
                    f"Warning: Judge selected out-of-bounds index: {selected_index}. Num unique candidates: {num_unique_candidates}")
            return None, reasoning + f" (Error: Index {selected_index} out of bounds for {num_unique_candidates} unique candidates)"
        return selected_index - 1, reasoning  # Convert to 0-based index for the unique list
    except json.JSONDecodeError as e:
        print(f"Warning: Judge response was not valid JSON: {judge_response_content}. Error: {e}")
        # Basic fallback: try to extract a number if JSON fails
        import re
        match_num = re.search(r'\b(\d+)\b', judge_response_content)
        if match_num:
            try:
                idx = int(match_num.group(1))
                if 1 <= idx <= num_unique_candidates:
                    return idx - 1, "Fallback: Extracted index from text. Original reasoning may be lost. Response: " + judge_response_content
            except ValueError:
                pass
        return None, f"Error parsing judge response (JSONDecodeError): {e}. Original: {judge_response_content}"
    except Exception as e:
        print(f"An unexpected error occurred while parsing judge response: {e}")
        return None, f"Unexpected error parsing judge response: {e}. Original: {judge_response_content}"


# --- Helper function for Deduplication and Formatting ---
def _format_candidates_for_judge(candidates: List[str], candidate_type: str = "Candidate") -> Tuple[str, List[str]]:
    """
    Deduplicates candidates, counts frequencies, and formats them for the judge.
    Returns the formatted string and the list of unique candidates in order.
    """
    if not candidates:
        return "No candidates generated.", []

    counts = Counter(candidates)
    unique_candidates = sorted(counts.keys(), key=lambda x: candidates.index(x))  # Preserve first appearance order

    formatted_string = ""
    code_block_type = "sql" if candidate_type.lower().startswith("sql") else "python" if candidate_type.lower().startswith("python") else ""

    for i, candidate in enumerate(unique_candidates):
        frequency = counts[candidate]
        freq_note = f"(generated {frequency} time{'s' if frequency > 1 else ''})"

        if candidate_type == "Insights Text Candidate":
             formatted_string += f"{candidate_type} {i + 1} {freq_note}:\n---\n{candidate}\n---\n\n"
        elif code_block_type: 
            formatted_string += f"{candidate_type} {i + 1} {freq_note}:\n```{code_block_type}\n{candidate}\n```\n\n"
        else:
            formatted_string += f"{candidate_type} {i + 1} {freq_note}:\n```\n{candidate}\n```\n\n"

    return formatted_string.strip(), unique_candidates


def _get_first_candidate(candidates: List[str], error_msg: str) -> str:
    """Safely gets the first candidate or returns an error string."""
    if candidates:
        return candidates[0]
    else:
        print(f"Warning: Candidate list empty for {error_msg}");
        return f"-- ERROR: Candidate list empty ({error_msg})"


def search_tables_and_schemas(state: AgentState) -> AgentState:
    """Retrieves table schemas relevant to the user's question."""
    print("\n--- Searching Tables and Schemas ---")
    docs_retrieved = retriever.invoke(state["question"])
    print(f"Number of documents retrieved: {len(docs_retrieved)}")
    # Filter out potentially non-JSON content before parsing
    tables_metadata = []
    for doc in docs_retrieved:

        try:
            # Handle potential variations in metadata structure if needed
            content = json.loads(doc.page_content)
            # Basic check for required keys before adding
            if isinstance(content, dict) and all(k in content for k in ["project_id", "dataset_id", "table_id"]):
                tables_metadata.append(content)
            else:
                print(f"Warning: Skipping document with unexpected structure: {doc.page_content[:100]}...")

        except json.JSONDecodeError:
            print(f"Warning: Skipping document with non-JSON content: {doc.page_content[:100]}...")

    if not tables_metadata:
        print("Warning: No valid table metadata found from retriever.")
        state["database_schemas"] = "No relevant table schemas found."
        # Decide if workflow should stop here or proceed with user potentially providing schema
        # For now, proceed, but SQL writer might fail.
        return state

    schemas = []

    bq_client = bigquery.Client(project=settings.project_id)
    fetched_tables = set()  # Avoid fetching schema for the same table multiple times if retrieved docs overlap

    for table_metadata in tables_metadata:
        project_id = table_metadata.get("project_id")
        dataset_id = table_metadata.get("dataset_id")
        table_id = table_metadata.get("table_id")
        full_table_id = f"{project_id}.{dataset_id}.{table_id}"

        if full_table_id in fetched_tables:
            continue

        print(f"Fetching schema for: {full_table_id}")
        try:
            schema = bq_functions.get_table_schema(bq_client, project_id, dataset_id, table_id)
            if schema:
                # Prepend full table ID for clarity in the combined schema string
                schemas.append(f"Schema for table `{full_table_id}`:\n{schema}")
                fetched_tables.add(full_table_id)
            else:
                print(f"Warning: Could not fetch schema for {full_table_id}")
        except Exception as e:
            print(f"Error fetching schema for {full_table_id}: {e}")

    if not schemas:
        print("Warning: Could not retrieve schemas for any identified tables.")
        state["database_schemas"] = "Could not retrieve schemas for identified tables."
    else:
        state["database_schemas"] = "\n---------------\n".join(schemas)

    print(f"Database Schemas Found:\n{state['database_schemas']}")
    return state


def agent_sql_writer_node(state: AgentState) -> AgentState:
    """Generates N SQL query candidates and stores them all."""
    n_runs = state.get('sql_writer_n_runs', 1)  # Default to 1 if not set
    print(f"\n--- Generating {n_runs} SQL Query Candidates ---")

    # Instantiate LLM dynamically
    llm_config = state.get('llm_config_sql_writer', {})
    print(f"SQL Writer using LLM config: {llm_config}")
    llm_sql_writer_instance = ChatVertexAI(**llm_config)  # Pass parameters

    prompt_template = ChatPromptTemplate.from_messages(("system", prompts.system_prompt_agent_sql_writer))

    chain = prompt_template | llm_sql_writer_instance;

    all_candidate_queries = []
    total_gen_input_tokens = 0
    total_gen_output_tokens = 0

    for i in range(n_runs):
        print(f"Generating SQL candidate {i + 1}/{n_runs}...")
        try:
            llm_response: AIMessage = chain.invoke({
                "question"        : state["question"],
                "database_schemas": state["database_schemas"]
            })
            query = utils.extract_code_block(content=llm_response.content, language="sql")
            all_candidate_queries.append(
                    query if query else f"-- Error: Could not extract SQL from response:\n{llm_response.content}")

            if llm_response.usage_metadata:
                total_gen_input_tokens += llm_response.usage_metadata.get("input_tokens", 0)
                total_gen_output_tokens += llm_response.usage_metadata.get("output_tokens", 0)
        except Exception as e:
            print(f"  Error generating SQL candidate {i + 1}: {e}")
            all_candidate_queries.append(f"-- Error: Generation failed: {e}")

    # Store all generated queries (including duplicates)
    state["sql_writer_candidate_queries"] = all_candidate_queries
    state["sql_writer_generation_input_tokens"] = total_gen_input_tokens
    state["sql_writer_generation_output_tokens"] = total_gen_output_tokens
    # Clear the main query field, it will be populated after judging
    state["query"] = ""
    print(f"SQL Writer generated {len(all_candidate_queries)} candidates.")
    print(f"SQL Writer Total Generation Tokens: Input={total_gen_input_tokens}, Output={total_gen_output_tokens}")
    return state


def agent_sql_validator_node(state: AgentState) -> AgentState:
    """Validates distinct SQL queries, attempts fixes, stores results."""
    print("\n--- Validating Distinct SQL Query Candidates ---")
    candidates = state.get('sql_writer_candidate_queries', [])
    max_retries = state.get('max_num_retries_debug', 2)
    # Get distinct candidates while preserving order
    distinct_candidates = list(dict.fromkeys(candidates))

    bq_client = bigquery.Client(project=settings.project_id)
    validation_results = []
    passed_validation_queries = []

    # Initialize total fixer tokens for this node's execution
    total_validator_input_tokens = 0
    total_validator_output_tokens = 0

    print(f"Attempting to validate {len(distinct_candidates)} distinct candidates...")

    # Instantiate LLM dynamically
    llm_config = state.get('llm_config_sql_validator', {})
    print(f"SQL Writer using LLM config: {llm_config}")
    llm_sql_validator_instance = ChatVertexAI(**llm_config)  # Pass parameters

    # Prepare the fixer chain (used inside the loop)
    fixer_prompt_template = ChatPromptTemplate.from_messages(
            ("system", prompts.system_prompt_agent_sql_validator_node))
    fixer_chain = fixer_prompt_template | llm_sql_validator_instance

    for i, original_query in enumerate(distinct_candidates):
        print(f"\nValidating Distinct Candidate {i + 1}:")
        print(f"Original Query:\n{original_query}")

        current_query = original_query
        passed = False
        final_error = None
        num_attempts = 0
        # sql_fixer_input_tokens = 0 # Track tokens per candidate if needed
        # sql_fixer_output_tokens = 0

        # --- Retry/Fix Loop for this distinct candidate ---
        while num_attempts <= max_retries:
            num_attempts += 1
            print(f"Attempt {num_attempts}/{max_retries + 1}...")

            # Basic check for invalid query text
            if not current_query or current_query.startswith("-- Error"):
                final_error = "Invalid or missing query text."
                print(f"  Result: FAILED ({final_error})")
                break  # No point retrying invalid text

            try:
                # Configure a dry run job
                job_config = bigquery.QueryJobConfig(dry_run=True, use_query_cache=False)
                bq_client.query(current_query, job_config=job_config)
                passed = True
                final_error = None
                print(f"  Result: PASSED (Dry Run)")
                break  # Exit loop on success

            except Exception as e:
                final_error = str(e)[:max_characters_error_msg_debug]
                print(f"  Result: FAILED (Dry Run). Error: {final_error}")

                # If validation failed and retries remain, attempt to fix
                if num_attempts <= max_retries:
                    print(f"  Attempting to fix query...")
                    try:
                        fixer_response: AIMessage = fixer_chain.invoke({
                            "query"          : current_query,
                            "error_msg_debug": final_error,
                            # Include schema if fixer prompt needs it
                            # "database_schemas": state["database_schemas"]
                        })
                        fixed_query = utils.extract_code_block(content=fixer_response.content, language="sql")

                        # --- Add token counting for the fixer LLM call ---
                        if fixer_response.usage_metadata:
                            total_validator_input_tokens += fixer_response.usage_metadata.get("input_tokens", 0)
                            total_validator_output_tokens += fixer_response.usage_metadata.get("output_tokens",
                                                                                               0)
                            print(
                                    f"  Fixer LLM Tokens (Attempt {num_attempts}): Input={input_tokens}, Output={output_tokens}")
                        # --- End token counting ---

                        if fixed_query and fixed_query != current_query:
                            print(f"  Query adjusted by LLM:\n{fixed_query}")
                            current_query = fixed_query  # Use fixed query for next attempt
                        else:
                            print("  Fixer LLM did not provide a new fix. Retrying with original query failed.")
                            # Break if fixer doesn't help? Or let it try again? Let's break.
                            break
                    except Exception as llm_e:
                        print(f"  Error invoking fixer LLM: {llm_e}")
                        # Break if fixer itself fails
                        break
                else:
                    # This was the last attempt, break the loop
                    break

        # --- Store results for this distinct candidate ---
        result_entry = {
            'original_query': original_query,
            'final_query'   : current_query if passed else original_query,
            # Store the version that passed, or original if failed
            'passed'        : passed,
            'error'         : final_error,
            'attempts'      : num_attempts
        }
        validation_results.append(result_entry)

        if passed:
            passed_validation_queries.append(current_query)  # Add the query that passed

    state['sql_validation_results'] = validation_results
    state['sql_passed_validation_queries'] = passed_validation_queries
    state['sql_validator_input_tokens'] = total_validator_input_tokens
    state['sql_validator_output_tokens'] = total_validator_output_tokens
    print(f"\nSQL Validation complete. {len(passed_validation_queries)} distinct queries passed.")
    return state


def agent_sql_judge_node(state: AgentState) -> AgentState:
    """Judges the SQL queries that passed validation."""
    print("\n--- Judging Passed SQL Queries ---")
    passed_queries = state.get('sql_passed_validation_queries', [])
    num_passed = len(passed_queries)

    print(f"Number of queries passed validation: {num_passed}")

    # Reset judge state fields
    state['query'] = ""
    state['sql_writer_judge_reasoning'] = ""
    state['sql_writer_judge_input_tokens'] = 0
    state['sql_writer_judge_output_tokens'] = 0

    if num_passed == 0:
        print("No queries passed validation. Cannot select a query.")
        state['query'] = "-- ERROR: No queries passed validation"
        state['sql_writer_judge_reasoning'] = "No queries passed validation."
        return state  # Proceed to execution node, which should handle the error

    elif num_passed == 1:
        selected_query = passed_queries[0]
        print("Only one query passed validation. Skipping judge LLM call.")
        state['query'] = selected_query
        state['sql_writer_judge_reasoning'] = "Skipped judge: Only one query passed validation."
        print(f"Selected query (single pass): {selected_query}")
        return state

    else:
        # More than one query passed, proceed with judging
        print(f"Judging {num_passed} passed queries...")

        # Format *only the passed* queries for the judge.
        # Using the helper, assuming passed_queries is the list of unique strings to judge.
        formatted_passed_candidates, unique_passed_queries = _format_candidates_for_judge(
                passed_queries, candidate_type="Passed SQL Query"
        )

        # Assumes prompt exists in prompts.py - RENAME PROMPT if needed
        # Make sure this prompt asks judge to choose from the *passed* queries.
        judge_system_prompt = prompts.system_prompt_agent_sql_judge_node
        judge_prompt_template = ChatPromptTemplate.from_messages([
            ("system", judge_system_prompt),
            ("human",
             "User Question: {question}\n\nDatabase Schemas:\n{database_schemas}\n\nPassed SQL Queries:\n{formatted_candidates}")
        ])
        judge_chain = judge_prompt_template | llm_sql_judge

        try:
            judge_response: AIMessage = judge_chain.invoke({
                "question"            : state["question"],
                "database_schemas"    : state["database_schemas"],
                "formatted_candidates": formatted_passed_candidates
            })
            print(f"SQL Judge Raw Response:\n{judge_response.content}")
            # Parse response based on the number of *unique passed* queries
            selected_index_0_based, reasoning = parse_judge_response(judge_response.content, len(unique_passed_queries))

            state["sql_writer_judge_reasoning"] = reasoning if reasoning else "Could not parse judge reasoning."

            if selected_index_0_based is not None and 0 <= selected_index_0_based < len(unique_passed_queries):
                state["query"] = unique_passed_queries[selected_index_0_based]  # Store selected QUERY STRING
                print(f"SQL Judge selected Candidate {selected_index_0_based + 1} from passed list.")
            else:
                print(f"Warning: SQL Judge failed selection/parsing. Defaulting to first passed query.")
                state["query"] = unique_passed_queries[0]
                state["sql_writer_judge_reasoning"] += " (Fallback: Used first passed query)"

            # Store token usage
            if judge_response.usage_metadata:
                state["sql_writer_judge_input_tokens"] = judge_response.usage_metadata.get("input_tokens", 0)
                state["sql_writer_judge_output_tokens"] = judge_response.usage_metadata.get("output_tokens", 0)
                print(
                        f"SQL Judge Tokens: Input={state['sql_writer_judge_input_tokens']}, Output={state['sql_writer_judge_output_tokens']}")

        except Exception as e:
            print(f"Error during SQL judging: {e}")
            # Fallback to first passed query on error
            state["query"] = unique_passed_queries[0] if unique_passed_queries else "-- ERROR: Judging failed"
            state["sql_writer_judge_reasoning"] = f"Error during judging: {e}. (Fallback: Used first passed query)"

        print(f"Selected SQL Query (after judge):\n{state['query']}")
        print(f"SQL Judge Reasoning: {state['sql_writer_judge_reasoning']}")
        return state


def execute_selected_sql_node(state: AgentState) -> AgentState:
    """Executes the single SQL query selected by the judge to get the DataFrame."""
    print("\n--- Executing Selected SQL Query ---")
    query = state.get('query')
    state['sql_execution_error'] = None  # Reset error
    state['df'] = pd.DataFrame()  # Reset df

    if not query or query.startswith("-- ERROR"):
        error_msg = "No valid SQL query available to execute."
        print(f"Error: {error_msg}")
        state['sql_execution_error'] = error_msg
        return state

    print(f"Executing query:\n{query}")
    try:
        bq_client = bigquery.Client(project=settings.project_id)
        df = bq_client.query(query).to_dataframe()
        state['df'] = df
        print(f"Query executed successfully. DataFrame shape: {df.shape}")
        if df.empty:
            print("Warning: Query executed successfully but returned an empty DataFrame.")
    except Exception as e:
        error_msg = f"Error executing selected SQL query: {str(e)}"
        print(error_msg)
        state['sql_execution_error'] = error_msg[:max_characters_error_msg_debug]  # Store truncated error

    return state


def agent_bi_expert_node(state: AgentState) -> AgentState:
    """Generates N BI visualization request candidates and stores them all."""
    n_runs = state.get('bi_expert_n_runs', 1)  # Default to 1 if not set
    print(f"\n--- Generating {n_runs} BI Visualization Request Candidates ---")
    # Instantiate LLM dynamically
    llm_config = state.get('llm_config_bi_expert', {})
    print(f"SQL Writer using LLM config: {llm_config}")
    llm_bi_expert_instance = ChatVertexAI(**llm_config)  # Pass parameters

    prompt_template = ChatPromptTemplate(("system", prompts.system_prompt_agent_bi_expert_node))

    chain = prompt_template | llm_bi_expert_instance

    all_candidate_requests = []
    total_gen_input_tokens = 0
    total_gen_output_tokens = 0

    # Prepare context once
    try:
        df_structure = str(state["df"].dtypes);
        df_sample = state["df"].head(5).to_string()
    except Exception as e:
        print(f"Warning: Could not get DataFrame info for BI Expert: {e}");
        df_structure = "Error";
        df_sample = "Error"

    for i in range(n_runs):
        print(f"Generating BI request candidate {i + 1}/{n_runs}...")
        try:
            llm_response: AIMessage = chain.invoke({
                "question"    : state["question"],
                "query"       : state["query"],  # The single, executed query
                "df_structure": df_structure,
                "df_sample"   : df_sample
            })
            request_text = llm_response.content.strip()
            all_candidate_requests.append(request_text if request_text else "-- Error: Empty response")

            if llm_response.usage_metadata:
                total_gen_input_tokens += llm_response.usage_metadata.get("input_tokens", 0)
                total_gen_output_tokens += llm_response.usage_metadata.get("output_tokens", 0)
        except Exception as e:
            print(f"  Error generating BI request candidate {i + 1}: {e}")
            all_candidate_requests.append(f"-- Error: Generation failed: {e}")

    # Store all generated requests
    state["bi_expert_candidate_requests"] = all_candidate_requests
    state["bi_expert_generation_input_tokens"] = total_gen_input_tokens
    state["bi_expert_generation_output_tokens"] = total_gen_output_tokens
    # Clear the main request field, it will be populated after judging (or skipping judge)
    state["visualization_request"] = ""

    all_requests = state.get('bi_expert_candidate_requests', [])
    # Deduplicate to check how many unique options we have
    unique_requests = list(dict.fromkeys(all_requests))
    num_unique = len(unique_requests)

    print(f"Generated {len(all_requests)} total BI requests, {num_unique} unique.")

    if num_unique <= 1:
        # Includes case where 0 requests were generated or n_runs was 1
        print("<= 1 unique BI request generated. Skipping BI Judge node.")
        # Select the first unique candidate (handles 0 or 1 unique)
        state['visualization_request'] = _get_first_candidate(unique_requests, "BI Expert")
        state['bi_expert_judge_reasoning'] = "Skipped judge: <= 1 unique request generated."
        state['bi_expert_judge_input_tokens'] = 0  # No judge tokens used
        state['bi_expert_judge_output_tokens'] = 0
        print(f"Selected BI request: {state['visualization_request']}")

    print(f"BI Expert generated {len(all_candidate_requests)} candidates.")
    print(f"BI Expert Total Generation Tokens: Input={total_gen_input_tokens}, Output={total_gen_output_tokens}")
    return state


# --- BI Judge Node ---
def agent_bi_judge_node(state: AgentState) -> AgentState:
    """Judges the distinct BI requests generated."""
    print("\n--- Judging Distinct BI Requests ---")
    all_requests = state.get('bi_expert_candidate_requests', [])

    # Deduplicate and format distinct requests
    formatted_unique_requests, unique_requests = _format_candidates_for_judge(
            all_requests, candidate_type="BI Request Candidate"
    )
    num_unique = len(unique_requests)

    print(f"Number of unique BI requests generated: {num_unique}")

    # Reset judge state fields
    state['visualization_request'] = ""
    state['bi_expert_judge_reasoning'] = ""
    state['bi_expert_judge_input_tokens'] = 0
    state['bi_expert_judge_output_tokens'] = 0

    if num_unique == 0:
        print("No unique BI requests were generated. Cannot select.")
        state['visualization_request'] = "-- ERROR: No BI requests generated"
        state['bi_expert_judge_reasoning'] = "No requests generated."
        return state  # Proceed to Python Gen, which should handle the error

    elif num_unique == 1:
        selected_request = unique_requests[0]
        print("Only one unique BI request generated. Skipping judge LLM call.")
        state['visualization_request'] = selected_request
        state['bi_expert_judge_reasoning'] = "Skipped judge: Only one unique request generated."
        print(f"Selected request (single unique): {selected_request}")
        return state

    else:
        # More than one unique request, proceed with judging
        print(f"Judging {num_unique} unique BI requests...")

        # Prepare context
        try:
            df_structure = str(state["df"].dtypes);
            df_sample = state["df"].head(5).to_string()
        except Exception as e:
            print(
                    f"Warning: Could not get DataFrame info for BI Judge: {e}");
            df_structure = "Error";
            df_sample = "Error"
        # Instantiate LLM dynamically
        llm_config = state.get('llm_config_bi_judge', {})
        print(f"SQL Writer using LLM config: {llm_config}")
        llm_bi_judge_instance = ChatVertexAI(**llm_config)

        # Assumes prompt exists in prompts.py
        judge_system_prompt = prompts.system_prompt_agent_bi_judge_node  # Make sure this exists!
        judge_prompt_template = ChatPromptTemplate.from_messages([
            ("system", judge_system_prompt),
            ("human",
             "User Question: {question}\nSQL Query:\n```sql\n{query}\n```\nData Summary:\nStructure & Data Types:\n{df_structure}\nSample Data:\n{df_sample}\n\nDistinct Visualization Requests:\n{formatted_candidates}")
        ])
        judge_chain = judge_prompt_template | llm_bi_judge_instance

        try:
            judge_response: AIMessage = judge_chain.invoke({
                "question"            : state["question"],
                "query"               : state["query"],
                "df_structure"        : df_structure,
                "df_sample"           : df_sample,
                "formatted_candidates": formatted_unique_requests  # Pass the formatted string
            })
            print(f"BI Judge Raw Response:\n{judge_response.content}")
            # Parse response based on the number of *unique* requests judged
            selected_index_0_based, reasoning = parse_judge_response(judge_response.content, len(unique_requests))

            state["bi_expert_judge_reasoning"] = reasoning if reasoning else "Could not parse judge reasoning."

            if selected_index_0_based is not None and 0 <= selected_index_0_based < len(unique_requests):
                state["visualization_request"] = unique_requests[
                    selected_index_0_based]  # Store selected REQUEST STRING
                print(f"BI Judge selected Candidate {selected_index_0_based + 1}.")
            else:
                print(f"Warning: BI Judge failed selection/parsing. Defaulting to first unique request.")
                state["visualization_request"] = unique_requests[0]
                state["bi_expert_judge_reasoning"] += " (Fallback: Used first unique request)"

            # Store token usage
            if judge_response.usage_metadata:
                state["bi_expert_judge_input_tokens"] = judge_response.usage_metadata.get("input_tokens", 0)
                state["bi_expert_judge_output_tokens"] = judge_response.usage_metadata.get("output_tokens", 0)
                print(
                        f"BI Judge Tokens: Input={state['bi_expert_judge_input_tokens']}, Output={state['bi_expert_judge_output_tokens']}")

        except Exception as e:
            print(f"Error during BI judging: {e}")
            # Fallback to first unique request on error
            state["visualization_request"] = unique_requests[0] if unique_requests else "-- ERROR: Judging failed"
            state["bi_expert_judge_reasoning"] = f"Error during judging: {e}. (Fallback: Used first unique request)"

        print(f"Selected Visualization Request (after judge):\n{state['visualization_request']}")
        print(f"BI Judge Reasoning: {state['bi_expert_judge_reasoning']}")
        return state


def agent_python_code_data_visualization_generator_node(state: AgentState) -> AgentState:
    """Generates N Python code candidates and stores them all."""
    n_runs = state.get('py_gen_n_runs', 1)  # Default to 1
    print(f"\n--- Generating {n_runs} Python Code Candidates ---")

    # Ensure visualization_request is valid before proceeding
    visualization_request = state.get("visualization_request")
    print('visualization_request: ', visualization_request)
    if not visualization_request or visualization_request.startswith("-- ERROR"):
        print("Error: Cannot generate Python code without a valid visualization request.")
        state[
            "py_gen_candidate_codes"] = ["# ERROR: Missing or invalid visualization request"] * n_runs  # Populate with error for each run
        state["python_code_data_visualization"] = "# ERROR: Missing or invalid visualization request"
        state["py_gen_generation_input_tokens"] = 0
        state["py_gen_generation_output_tokens"] = 0
        return state

    # Instantiate LLM dynamically
    llm_config = state.get('llm_config_python_writer', {})
    print(f"SQL Writer using LLM config: {llm_config}")
    llm_python_writer_instance = ChatVertexAI(**llm_config)

    prompt_template = ChatPromptTemplate.from_messages(
            ("system", prompts.system_prompt_agent_python_code_data_visualization_generator_node))
    chain = prompt_template | llm_python_writer_instance

    all_candidate_codes = []
    total_gen_input_tokens = 0
    total_gen_output_tokens = 0

    # Prepare context once
    try:
        df_structure = str(state["df"].dtypes);
        df_sample = state["df"].head(5).to_string()
    except Exception as e:
        print(f"Warning: Could not get DataFrame info for Python Generator: {e}");
        df_structure = "Error";
        df_sample = "Error"

    for i in range(n_runs):
        print(f"Generating Python code candidate {i + 1}/{n_runs}...")
        try:
            llm_response: AIMessage = chain.invoke({
                "visualization_request": visualization_request,
                "df_structure"         : df_structure,
                "df_sample"            : df_sample
            })
            code = utils.extract_code_block(content=llm_response.content, language="python")
            all_candidate_codes.append(
                    code if code else f"# ERROR: Could not extract Python code from response:\n{llm_response.content}")

            if llm_response.usage_metadata:
                total_gen_input_tokens += llm_response.usage_metadata.get("input_tokens", 0)
                total_gen_output_tokens += llm_response.usage_metadata.get("output_tokens", 0)
        except Exception as e:
            print(f"  Error generating Python code candidate {i + 1}: {e}")
            all_candidate_codes.append(f"# ERROR: Generation failed: {e}")

    # Store all generated codes
    state["py_gen_candidate_codes"] = all_candidate_codes
    state["py_gen_generation_input_tokens"] = total_gen_input_tokens
    state["py_gen_generation_output_tokens"] = total_gen_output_tokens
    # Clear the main code field, it will be populated after judging (or skipping judge)
    state["python_code_data_visualization"] = ""

    all_codes = state.get('py_gen_candidate_codes', [])
    # Deduplicate to check how many unique options we have
    unique_codes = list(dict.fromkeys(all_codes))
    num_unique = len(unique_codes)

    print(f"Generated {len(all_codes)} total Python codes, {num_unique} unique.")

    if num_unique <= 1:
        # Select the first unique candidate (handles 0 or 1 unique)
        state['python_code_data_visualization'] = _get_first_candidate(unique_codes, "Python Generator")
        state['py_gen_judge_reasoning'] = "Skipped judge: <= 1 unique code generated."
        state['py_gen_judge_input_tokens'] = 0  # No judge tokens used
        state['py_gen_judge_output_tokens'] = 0
        print(f"Selected Python code:\n```python\n{state['python_code_data_visualization']}\n```")

    print(f"Python Generator generated {len(all_candidate_codes)} candidates.")
    print(f"Python Generator Total Generation Tokens: Input={total_gen_input_tokens}, Output={total_gen_output_tokens}")
    return state


def agent_python_judge_node(state: AgentState) -> AgentState:
    """Judges the distinct Python code snippets generated."""
    print("\n--- Judging Distinct Python Code Candidates ---")
    all_codes = state.get('py_gen_candidate_codes', [])

    # Deduplicate and format distinct codes
    formatted_unique_codes, unique_codes = _format_candidates_for_judge(
            all_codes, candidate_type="Python Code Candidate"  # Ensure specific type for formatting
    )
    num_unique = len(unique_codes)

    print(f"Number of unique Python codes generated: {num_unique}")

    # Reset judge state fields
    state['python_code_data_visualization'] = ""
    state['py_gen_judge_reasoning'] = ""
    state['py_gen_judge_input_tokens'] = 0
    state['py_gen_judge_output_tokens'] = 0

    if num_unique == 0:
        print("No unique Python codes were generated. Cannot select.")
        state['python_code_data_visualization'] = "# ERROR: No Python codes generated"
        state['py_gen_judge_reasoning'] = "No codes generated."
        return state  # Proceed to Validator, which should handle the error

    elif num_unique == 1:
        selected_code = unique_codes[0]
        print("Only one unique Python code generated. Skipping judge LLM call.")
        state['python_code_data_visualization'] = selected_code
        state['py_gen_judge_reasoning'] = "Skipped judge: Only one unique code generated."
        print(f"Selected code (single unique):\n```python\n{selected_code}\n```")
        return state

    else:
        # More than one unique code, proceed with judging
        print(f"Judging {num_unique} unique Python codes...")

        # Prepare context
        try:
            df_structure = str(state["df"].dtypes);
            df_sample = state["df"].head(5).to_string()
        except Exception as e:
            print(
                    f"Warning: Could not get DataFrame info for Python Judge: {e}");
            df_structure = "Error";
            df_sample = "Error"

        # Instantiate LLM dynamically
        llm_config = state.get('llm_config_python_judge', {})
        print(f"SQL Writer using LLM config: {llm_config}")
        llm_python_judge_instance = ChatVertexAI(**llm_config)

        # Assumes prompt exists in prompts.py
        judge_system_prompt = prompts.system_prompt_agent_python_judge_node  # Make sure this exists!
        judge_prompt_template = ChatPromptTemplate.from_messages([
            ("system", judge_system_prompt),
            ("human",
             "Visualization Request:\n{visualization_request}\n\nData Summary:\nStructure & Data Types:\n{df_structure}\nSample Data:\n{df_sample}\n\nDistinct Python Code Candidates:\n{formatted_candidates}")
        ])
        judge_chain = judge_prompt_template | llm_python_judge_instance

        try:
            judge_response: AIMessage = judge_chain.invoke({
                "visualization_request": state["visualization_request"],
                "df_structure"         : df_structure,
                "df_sample"            : df_sample,
                "formatted_candidates" : formatted_unique_codes  # Pass the formatted string
            })
            print(f"Python Judge Raw Response:\n{judge_response.content}")
            # Parse response based on the number of *unique* codes judged
            selected_index_0_based, reasoning = parse_judge_response(judge_response.content, len(unique_codes))

            state["py_gen_judge_reasoning"] = reasoning if reasoning else "Could not parse judge reasoning."

            if selected_index_0_based is not None and 0 <= selected_index_0_based < len(unique_codes):
                state["python_code_data_visualization"] = unique_codes[
                    selected_index_0_based]  # Store selected CODE STRING
                print(f"Python Judge selected Candidate {selected_index_0_based + 1}.")
            else:
                print(f"Warning: Python Judge failed selection/parsing. Defaulting to first unique code.")
                state["python_code_data_visualization"] = unique_codes[0]
                state["py_gen_judge_reasoning"] += " (Fallback: Used first unique code)"

            # Store token usage
            if judge_response.usage_metadata:
                state["py_gen_judge_input_tokens"] = judge_response.usage_metadata.get("input_tokens", 0)
                state["py_gen_judge_output_tokens"] = judge_response.usage_metadata.get("output_tokens", 0)
                print(
                        f"Python Judge Tokens: Input={state['py_gen_judge_input_tokens']}, Output={state['py_gen_judge_output_tokens']}")

        except Exception as e:
            print(f"Error during Python code judging: {e}")
            # Fallback to first unique code on error
            state["python_code_data_visualization"] = unique_codes[0] if unique_codes else "# ERROR: Judging failed"
            state["py_gen_judge_reasoning"] = f"Error during judging: {e}. (Fallback: Used first unique code)"

        print(f"Selected Python Code (after judge):\n```python\n{state['python_code_data_visualization']}\n```")
        print(f"Python Judge Reasoning: {state['py_gen_judge_reasoning']}")
        return state


def agent_python_code_data_visualization_validator_node(state: AgentState) -> AgentState:
    print("\n\n### Validating data visualization code:")

    try:
        df = state["df"]
        # Create a dictionary to store the executed variables for the python code generated

        exec_globals = {"df": df, "pd": pd}  # Added pd for safety, if generated code uses it

        exec(state["python_code_data_visualization"], exec_globals)

        state["python_code_store_variables_dict"] = exec_globals
        state["result_debug_python_code_data_visualization"] = "Pass"
        state["error_msg_debug_python_code_data_visualization"] = ""

        print(f"result: {state['result_debug_python_code_data_visualization']}")

        return state

    except Exception as e:
        state["num_retries_debug_python_code_data_visualization"] += 1

        # return False, f"Error validating query: {str(e)}"
        state["result_debug_python_code_data_visualization"] = "Not Pass"
        state["error_msg_debug_python_code_data_visualization"] = str(e)[0:max_characters_error_msg_debug]
        print(f"result: {state['result_debug_python_code_data_visualization']}")
        print(f'error message: {state["error_msg_debug_python_code_data_visualization"]}')

        if state['num_retries_debug_python_code_data_visualization'] <= state['max_num_retries_debug']:
            print("\n### Trying to fix the plotly code:")

            # Instantiate LLM dynamically
            llm_config = state.get('llm_config_python_validator', {})
            print(f"SQL Writer using LLM config: {llm_config}")
            llm_python_validator_instance = ChatVertexAI(**llm_config)

            prompt_template = ChatPromptTemplate.from_messages(
                    [("system", prompts.system_prompt_agent_python_code_data_visualization_validator_node)])

            chain = prompt_template | llm_python_validator_instance

            llm_response: AIMessage = chain.invoke({
                "python_code_data_visualization": state["python_code_data_visualization"],
                "error_msg_debug"               : state["error_msg_debug_python_code_data_visualization"]
            })

            state["python_code_data_visualization"] = utils.extract_code_block(content=llm_response.content,
                                                                               language="python")

            print(f"\n### Plotly code adjusted:\n {state['python_code_data_visualization']}")

            if llm_response.usage_metadata:
                state["py_fixer_input_tokens"] += llm_response.usage_metadata.get("input_tokens", 0)
                state["py_fixer_output_tokens"] += llm_response.usage_metadata.get("output_tokens", 0)

                print(
                        f"Python Code Fixer Tokens (Attempt {state['num_retries_debug_python_code_data_visualization']}): Input={llm_response.usage_metadata.get('input_tokens', 0)}, Output={llm_response.usage_metadata.get('output_tokens', 0)}")

        else:
            print(f"Max retries for Python code validation reached ({state['max_num_retries_debug']}).")

        return state


def agent_insights_generator_node(state: AgentState) -> AgentState:
    n_runs = state.get('insights_n_runs', 1) # Get n_runs for insights
    print(f"\n--- Generating {n_runs} Insights Candidate(s) ---")

    # Initialize fields for multiple runs
    all_candidate_insights = []
    total_gen_input_tokens = 0
    total_gen_output_tokens = 0
    
    # Prepare common inputs once
    current_df_for_insights = state.get('df')
    df_viz_summary_str = "DataFrame is empty or not available."
    if current_df_for_insights is not None and not current_df_for_insights.empty:
        buffer = io.StringIO()
        buffer.write("DataFrame Head:\n" + current_df_for_insights.head().to_string() + "\n\nDataFrame Info:\n")
        current_df_for_insights.info(buf=buffer)
        df_viz_summary_str = buffer.getvalue().strip()
    elif current_df_for_insights is not None and current_df_for_insights.empty:
        df_viz_summary_str = "DataFrame `df` (used for visualization) is empty."

    fig_object = state.get('python_code_store_variables_dict', {}).get('fig')
    plot_json_str = "{}"
    if fig_object:
        try: plot_json_str = plotly.io.to_json(fig_object, pretty=False).strip()
        except Exception as e: plot_json_str = f'{{"error": "Could not convert plot to JSON: {str(e)}"}}'
    else: plot_json_str = '{"message": "Plot (variable `fig`) not available or not generated."}'
    
    mtss_context_str = state.get('insights_context_mtss_docs', "MTSS context not provided in state.").strip()
    user_stories_str = state.get('insights_context_user_stories', "User stories context not provided in state.").strip()
    correlations_str = state.get('insights_context_correlations', "Correlation summary not provided in state.").strip()
    
    llm_config = state.get('llm_config_insights_generator', {})
    print(f"Insights Generator using LLM config: {llm_config}")
    insights_llm_instance = ChatVertexAI(**llm_config)
    
    prompt_input_dict = {
        "df_viz_summary": df_viz_summary_str if df_viz_summary_str else "Data summary not available.",
        "plot_json": plot_json_str if plot_json_str else '{"message": "Plot JSON not available."}',
        "mtss_context": mtss_context_str if mtss_context_str else "MTSS context not available.",
        "user_stories_content": user_stories_str if user_stories_str else "User stories not available.",
        "strong_corrs_summary": correlations_str if correlations_str else "Correlation summary not available."
    }
    
    # Check for empty inputs before forming the prompt template
    for key, value in prompt_input_dict.items():
        if not isinstance(value, str) or not value.strip():
            print(f"CRITICAL WARNING: Input for placeholder '{key}' is empty or only whitespace for Insights LLM! Value: '{value}'")
            prompt_input_dict[key] = f"<{key} content was empty or not a string>"

    raw_prompt_template_str = prompts.system_prompt_agent_insights_generator_node
    if not raw_prompt_template_str or not raw_prompt_template_str.strip():
        state['generated_insights'] = "Error: System prompt template for insights is empty."
        state['insights_candidate_texts'] = [state['generated_insights']] * n_runs
        return state

    prompt_template = ChatPromptTemplate.from_messages(("system", raw_prompt_template_str))
    chain = prompt_template | insights_llm_instance

    for i in range(n_runs):
        print(f"Generating Insights Candidate {i + 1}/{n_runs}...")
        try:
            llm_response: AIMessage = chain.invoke(prompt_input_dict)
            insight_text = llm_response.content.strip()
            all_candidate_insights.append(insight_text if insight_text else "-- Error: Empty insight response")
            if llm_response.usage_metadata:
                total_gen_input_tokens += llm_response.usage_metadata.get("input_tokens", 0)
                total_gen_output_tokens += llm_response.usage_metadata.get("output_tokens", 0)
        except Exception as e:
            print(f"  Error generating insight candidate {i + 1}: {e}")
            all_candidate_insights.append(f"-- Error: Insight Generation failed: {e}")

    state['insights_candidate_texts'] = all_candidate_insights
    state['insights_generator_input_tokens'] = total_gen_input_tokens
    state['insights_generator_output_tokens'] = total_gen_output_tokens
    state['generated_insights'] = "" # Clear, to be set by judge or this node if skipping judge

    valid_unique_insights = list(dict.fromkeys(text for text in all_candidate_insights if text and not text.startswith("-- Error")))
    if len(valid_unique_insights) <= 1:
        state['generated_insights'] = _get_first_candidate(valid_unique_insights, "Insights Generator", error_indicator="-- Error:")
        state['insights_judge_reasoning'] = "Skipped judge: <= 1 unique valid insight generated."
        state['insights_judge_input_tokens'] = 0
        state['insights_judge_output_tokens'] = 0
        print(f"Selected Insight (skip judge): {state['generated_insights'][:200]}...")
    
    print(f"Insights Generator produced {len(all_candidate_insights)} candidates.")
    print(f"Insights Generator Total Tokens: Input={total_gen_input_tokens}, Output={total_gen_output_tokens}")
    return state

def agent_insights_judge_node(state: AgentState) -> AgentState:
    print("\n--- Judging Distinct Insight Candidates ---")
    all_insight_candidates = state.get('insights_candidate_texts', [])
    
    # Filter out error strings before formatting for judge
    valid_insight_candidates = [text for text in all_insight_candidates if text and not text.startswith("-- Error")]
    
    formatted_candidates_str, unique_insights = _format_candidates_for_judge(
        valid_insight_candidates, candidate_type="Insights Text Candidate" # Use a descriptive type
    )
    num_unique = len(unique_insights)

    print(f"Number of unique insight texts to judge: {num_unique}")
    # Initialize judge fields
    state['generated_insights'] = "" # Will be set by the judge
    state['insights_judge_reasoning'] = ""
    state['insights_judge_input_tokens'] = 0
    state['insights_judge_output_tokens'] = 0

    if not unique_insights:
        state['generated_insights'] = "-- ERROR: No valid insights to judge"
        state['insights_judge_reasoning'] = "No valid insights were generated to be judged."
        return state
    # If only one unique insight, it should have been handled by skip logic in agent_insights_generator_node

    # Prepare context for the judge (similar to generator but without N-runs)
    current_df_for_insights = state.get('df')
    df_viz_summary_str = "DataFrame is empty or not available."
    if current_df_for_insights is not None and not current_df_for_insights.empty:
        buffer = io.StringIO()
        buffer.write("DataFrame Head:\n" + current_df_for_insights.head().to_string() + "\n\nDataFrame Info:\n")
        current_df_for_insights.info(buf=buffer)
        df_viz_summary_str = buffer.getvalue().strip()
    elif current_df_for_insights is not None and current_df_for_insights.empty:
         df_viz_summary_str = "DataFrame `df` (used for visualization) is empty."

    fig_object = state.get('python_code_store_variables_dict', {}).get('fig')
    plot_json_str = "{}"
    if fig_object:
        try: plot_json_str = plotly.io.to_json(fig_object, pretty=False).strip()
        except Exception as e: plot_json_str = f'{{"error": "Could not convert plot to JSON: {str(e)}"}}'
    else: plot_json_str = '{"message": "Plot (variable `fig`) not available or not generated."}'
    
    mtss_context_str = state.get('insights_context_mtss_docs', "MTSS context not provided.").strip()
    user_stories_str = state.get('insights_context_user_stories', "User stories context not provided.").strip()
    correlations_str = state.get('insights_context_correlations', "Correlation summary not provided.").strip()

    llm_config = state.get('llm_config_insights_judge', {})
    print(f"Insights Judge using LLM config: {llm_config}")
    insights_judge_llm_instance = ChatVertexAI(**llm_config)

    judge_prompt_template = ChatPromptTemplate.from_messages(
        ("system", prompts.system_prompt_agent_insights_judge_node)
    )
    judge_chain = judge_prompt_template | insights_judge_llm_instance

    prompt_input_dict = {
        "question": state["question"], # Pass original question to judge
        "df_viz_summary": df_viz_summary_str if df_viz_summary_str else "Data summary not available.",
        "plot_json": plot_json_str if plot_json_str else '{"message": "Plot JSON not available."}',
        "mtss_context": mtss_context_str if mtss_context_str else "MTSS context not available.",
        "user_stories_content": user_stories_str if user_stories_str else "User stories not available.",
        "strong_corrs_summary": correlations_str if correlations_str else "Correlation summary not available.",
        "formatted_candidates": formatted_candidates_str # The list of insights to judge
    }
    
    try:
        print("Invoking Insights Judge LLM...")
        judge_response: AIMessage = judge_chain.invoke(prompt_input_dict)
        print(f"Insights Judge Raw Response:\n{judge_response.content}")
        
        selected_idx, reasoning_text = parse_judge_response(judge_response.content, len(unique_insights))
        state["insights_judge_reasoning"] = reasoning_text if reasoning_text else "Could not parse insights judge reasoning."
        
        if selected_idx is not None and 0 <= selected_idx < len(unique_insights):
            state["generated_insights"] = unique_insights[selected_idx]
            print(f"Insights Judge selected Candidate {selected_idx + 1}.")
        else:
            state["generated_insights"] = unique_insights[0] # Fallback
            state["insights_judge_reasoning"] += " (Fallback: Used first unique insight)"
            print(f"Warning: Insights Judge failed selection/parsing. Defaulting to first unique insight.")

        if judge_response.usage_metadata:
            state["insights_judge_input_tokens"] = judge_response.usage_metadata.get("input_tokens", 0)
            state["insights_judge_output_tokens"] = judge_response.usage_metadata.get("output_tokens", 0)
            print(f"Insights Judge Tokens: Input={state['insights_judge_input_tokens']}, Output={state['insights_judge_output_tokens']}")

    except Exception as e:
        print(f"Error during Insights judging: {e}")
        state["generated_insights"] = unique_insights[0] if unique_insights else "-- ERROR: Insights Judging failed"
        state["insights_judge_reasoning"] = f"Error during judging: {e}. (Fallback: Used first unique insight)"
    
    print(f"Selected Insight (after judge): {state['generated_insights'][:200]}...")
    print(f"Insights Judge Reasoning: {state['insights_judge_reasoning']}")
    return state


# Conditional edge after SQL Judge
def decide_after_sql_judge(state: AgentState) -> str:
    """Routes to execution if judge succeeded, otherwise ends."""
    print(f"--- Decision: After SQL Judge ---")
    query = state.get('query', '')
    # Check if the selected query indicates an error (either no pass or judge fail)
    if not query or query.startswith("-- ERROR"):
        print("Ending: No valid SQL query was selected.")
        return END
    else:
        print("Proceeding to execute selected SQL.")
        return "execute_selected_sql_node"


# Conditional edge after SQL Execution
def decide_after_sql_execution(state: AgentState) -> str:
    """Routes to BI expert if execution succeeded, otherwise ends."""
    print(f"--- Decision: After SQL Execution ---")
    if state.get('sql_execution_error') is None and not state.get('df', pd.DataFrame()).empty:
        print("SQL Execution successful. Proceeding to BI Expert.")
        return "agent_bi_expert_node"
    elif state.get('sql_execution_error') is None and state.get('df', pd.DataFrame()).empty:
        print("Warning: SQL Execution returned empty DataFrame. Proceeding to BI Expert.")
        # Proceed even if df is empty, let BI/Python handle it
        return "agent_bi_expert_node"
    else:
        print(f"Ending: SQL Execution failed. Error: {state.get('sql_execution_error')}")
        return END


def decide_after_bi_expert(state: AgentState) -> str:
    """Decides whether to judge BI requests or use the single candidate."""
    print(f"\n--- Decision: After BI Expert ---")
    all_requests = state.get('bi_expert_candidate_requests', [])
    # Deduplicate to check how many unique options we have
    unique_requests = list(dict.fromkeys(all_requests))
    num_unique = len(unique_requests)

    print(f"Generated {len(all_requests)} total BI requests, {num_unique} unique.")

    if num_unique <= 1:
        # Proceed directly to Python generator
        print("Proceeding to Python generator node.")
        return "agent_python_code_data_visualization_generator_node"
    else:
        # More than 1 unique request, go to judge
        print("Proceeding to BI Judge node.")
        return "agent_bi_judge_node"


# Conditional edge after BI Judge (if it runs)
def decide_after_bi_judge(state: AgentState) -> str:
    """Routes to Python generator if judge succeeded, otherwise ends."""
    print(f"--- Decision: After BI Judge ---")
    request = state.get('visualization_request', '')
    # Check if the selected request indicates an error
    if not request or request.startswith("-- ERROR"):
        print("Ending: No valid BI request was selected.")
        return END
    else:
        print("Proceeding to Python Code Generation.")
        return "agent_python_code_data_visualization_generator_node"


# Conditional edge after Python Judge (if it runs)
def decide_after_python_judge(state: AgentState) -> str:
    """Routes to Python validator if judge succeeded, otherwise ends."""
    print(f"--- Decision: After Python Judge ---")
    code = state.get('python_code_data_visualization', '')
    # Check if the selected code indicates an error
    if not code or code.startswith("# ERROR"):
        print("Ending: No valid Python code was selected by the judge.")
        return END
    else:
        print("Proceeding to Python Code Validation.")
        # Reset Python validation retry counter before validation
        return "agent_python_code_data_visualization_validator_node"


def decide_after_python_generator(state: AgentState) -> str:
    """Decides whether to judge Python codes or use the single candidate."""
    print(f"\n--- Decision: After Python Code Generator ---")
    all_codes = state.get('py_gen_candidate_codes', [])
    # Deduplicate to check how many unique options we have
    unique_codes = list(dict.fromkeys(all_codes))
    num_unique = len(unique_codes)

    print(f"Generated {len(all_codes)} total Python codes, {num_unique} unique.")

    if num_unique <= 1:
        # Includes case where 0 codes were generated or n_runs was 1
        print("<= 1 unique Python code generated. Skipping Python Judge node.")
        # Proceed directly to Python validator
        return "agent_python_code_data_visualization_validator_node"
    else:
        # More than 1 unique code, go to judge
        print("Proceeding to Python Judge node.")
        return "agent_python_judge_node"


def decide_after_python_validation(state: AgentState) -> str:
    if state['result_debug_python_code_data_visualization'] == "Pass":
        return "agent_insights_generator_node"  # MODIFIED: Go to insights

    if state['num_retries_debug_python_code_data_visualization'] >= state['max_num_retries_debug']:
        print("Python code validation failed after max retries. Ending.")
        return END

    return 'agent_python_code_data_visualization_validator_node'

# --- Conditional Edge after Insights Generator ---
def decide_after_insights_generator(state: AgentState) -> str:
    """Decides whether to judge insights or end."""
    print(f"\n--- Decision: After Insights Generator ---")
    # If generated_insights is already set, it means judge was skipped (<=1 unique candidate)
    if state.get("generated_insights") and not state.get("generated_insights").startswith("-- Error:"):
        print("Insights already selected (<=1 unique or judge skipped). Ending.")
        return END
    elif state.get("generated_insights") and state.get("generated_insights").startswith("-- Error:"):
        print("Ending: Error in Insights generator or no valid insights selected by generator.")
        return END
    else: 
        # This means multiple unique candidates were generated and 'generated_insights' is still empty
        print("Multiple unique insights generated. Proceeding to Insights Judge node.")
        return "agent_insights_judge_node"


# --- Workflow Definition ---
workflow = StateGraph(state_schema=AgentState)

# --- Define nodes ---
workflow.add_node("search_tables_and_schemas", search_tables_and_schemas)
workflow.add_node("agent_sql_writer_node", agent_sql_writer_node)
workflow.add_node("agent_sql_validator_node", agent_sql_validator_node)
workflow.add_node("agent_sql_judge_node", agent_sql_judge_node)  # Judges passed
workflow.add_node("execute_selected_sql_node", execute_selected_sql_node)  # Executes final choice
workflow.add_node("agent_bi_expert_node", agent_bi_expert_node)
workflow.add_node("agent_bi_judge_node", agent_bi_judge_node)
workflow.add_node("agent_python_code_data_visualization_generator_node",
                  agent_python_code_data_visualization_generator_node)  # Runs n times
workflow.add_node("agent_python_code_data_visualization_validator_node",
                  agent_python_code_data_visualization_validator_node)
workflow.add_node("agent_python_judge_node", agent_python_judge_node)
workflow.add_node("agent_insights_generator_node", agent_insights_generator_node)
workflow.add_node("agent_insights_judge_node", agent_insights_judge_node)

# --- Define edges ---
workflow.set_entry_point("search_tables_and_schemas")
# SQL Path
workflow.add_edge("search_tables_and_schemas", "agent_sql_writer_node")
workflow.add_edge("agent_sql_writer_node", "agent_sql_validator_node")  # Writer -> Validate Distinct
workflow.add_edge("agent_sql_validator_node", "agent_sql_judge_node")  # Validate Distinct -> Judge Passed
workflow.add_conditional_edges(
        "agent_sql_judge_node",
        decide_after_sql_judge,
        {
            "execute_selected_sql_node": "execute_selected_sql_node",
            END                        : END
        }
)
workflow.add_conditional_edges(
        "execute_selected_sql_node",
        decide_after_sql_execution,
        {
            "agent_bi_expert_node": "agent_bi_expert_node",
            END                   : END
        }
)
# BI Path
workflow.add_conditional_edges(
        "agent_bi_expert_node",
        decide_after_bi_expert,
        {
            "agent_bi_judge_node"                                : "agent_bi_judge_node",  # If > 1 unique request
            "agent_python_code_data_visualization_generator_node": "agent_python_code_data_visualization_generator_node"
            # If <= 1 unique request
        }
)

workflow.add_conditional_edges(
        "agent_bi_judge_node",
        decide_after_bi_judge,
        {
            "agent_python_code_data_visualization_generator_node": "agent_python_code_data_visualization_generator_node",
            END                                                  : END
        }
)
# Python Path
workflow.add_conditional_edges(
        "agent_python_code_data_visualization_generator_node",
        decide_after_python_generator,
        {
            "agent_python_judge_node"                            : "agent_python_judge_node",  # If > 1 unique code
            "agent_python_code_data_visualization_validator_node": "agent_python_code_data_visualization_validator_node"
            # If <= 1 unique code
        }
)
workflow.add_conditional_edges(
        "agent_python_judge_node",
        decide_after_python_judge,
        {
            "agent_python_code_data_visualization_validator_node": "agent_python_code_data_visualization_validator_node",
            END                                                  : END
        }
)
workflow.add_conditional_edges('agent_python_code_data_visualization_validator_node', decide_after_python_validation, {'agent_insights_generator_node': 'agent_insights_generator_node', 'agent_python_code_data_visualization_validator_node': 'agent_python_code_data_visualization_validator_node', END: END})

# Edges for Insights Path
workflow.add_conditional_edges('agent_insights_generator_node', decide_after_insights_generator, {
    'agent_insights_judge_node': 'agent_insights_judge_node', # Go to judge if needed
    END: END # End if judge was skipped by generator
})
workflow.add_edge('agent_insights_judge_node', END) # Insights judge is the final step in this branch

app = workflow.compile()


### Run workflow

def run_workflow(question: str,
                 user_stories_text: str = "User stories not provided.",
                 mtss_docs_text: str = "MTSS documentation not provided.",
                 correlations_text: str = "Correlations summary not provided.",
                 **kwargs) -> dict:
    """
    Runs the complete LangGraph workflow with customizable initial state and LLM parameters.
    Args:
        question: The user's question (required).
        **kwargs: Arbitrary keyword arguments.
                  - To override initial state: e.g., sql_writer_n_runs=5
                  - To set LLM params: e.g., llm_sql_writer_model_name="gemini-1.0-pro",
                                           llm_sql_writer_temperature=0.7,
                                           llm_bi_expert_max_tokens=500
    Returns:
        The final state dictionary of the workflow.
    """
    print(f"--- Starting Workflow ---")
    print(f"Question: {question}")

    # Default LLM configurations
    default_llm_params = {
        "project" : settings.project_id,
        "location": "us-central1"
        # Add other truly global defaults if any, e.g., 'max_retries': 3
    }

    default_sql_writer_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.5, **default_llm_params}
    default_sql_validator_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.2, **default_llm_params}
    default_sql_judge_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.1, **default_llm_params}
    default_bi_expert_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.5, **default_llm_params}
    default_bi_judge_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.1, **default_llm_params}
    default_python_writer_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.3, **default_llm_params}
    default_python_validator_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.2,
                                       **default_llm_params}
    default_python_judge_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.1, **default_llm_params}
    default_insights_generator_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.4,
                                         **default_llm_params}
    default_insights_judge_config = {"model_name": "gemini-2.0-flash-lite-001", "temperature": 0.4,
                                         **default_llm_params}

    default_initial_state = AgentState(
            question=question,
            database_schemas="",
            max_num_retries_debug=2,

            # LLM Configs
            llm_config_sql_writer=default_sql_writer_config.copy(),
            llm_config_sql_validator=default_sql_validator_config.copy(),
            llm_config_sql_judge=default_sql_judge_config.copy(),
            llm_config_bi_expert=default_bi_expert_config.copy(),
            llm_config_bi_judge=default_bi_judge_config.copy(),
            llm_config_python_writer=default_python_writer_config.copy(),
            llm_config_python_validator=default_python_validator_config.copy(),
            llm_config_python_judge=default_python_judge_config.copy(),
            llm_config_insights_generator=default_insights_generator_config.copy(),
            llm_config_insights_judge=default_insights_judge_config.copy(),

            # SQL Writer
            sql_writer_n_runs=1,
            sql_writer_candidate_queries=[],
            sql_writer_generation_input_tokens=0,
            sql_writer_generation_output_tokens=0,

            # SQL Validator
            sql_validation_results=[],
            sql_passed_validation_queries=[],
            sql_validator_input_tokens=0,
            sql_validator_output_tokens=0,

            # SQL Judge
            query="",
            sql_writer_judge_reasoning=None,
            sql_writer_judge_input_tokens=0,
            sql_writer_judge_output_tokens=0,

            # SQL Executor
            df=pd.DataFrame(),
            sql_execution_error=None,

            # BI Expert
            visualization_request="",
            bi_expert_n_runs=1,
            bi_expert_candidate_requests=[],
            bi_expert_judge_reasoning=None,
            bi_expert_generation_input_tokens=0,
            bi_expert_generation_output_tokens=0,
            bi_expert_judge_input_tokens=0,
            bi_expert_judge_output_tokens=0,
            
            # Python Gen
            python_code_data_visualization="",
            py_gen_n_runs=1,
            py_gen_candidate_codes=[],
            py_gen_judge_reasoning=None,
            py_gen_generation_input_tokens=0,
            py_gen_generation_output_tokens=0,
            py_gen_judge_input_tokens=0,
            py_gen_judge_output_tokens=0,

            # Python Validator
            python_code_store_variables_dict={},
            num_retries_debug_python_code_data_visualization=0,
            result_debug_python_code_data_visualization="",
            error_msg_debug_python_code_data_visualization="",
            py_fixer_input_tokens=0,
            py_fixer_output_tokens=0,

            # Insights Agent
            insights_n_runs=2,
            insights_context_user_stories=user_stories_text,
            insights_context_mtss_docs=mtss_docs_text,
            insights_context_correlations=correlations_text,
            insights_candidate_texts=[],
            generated_insights="",
            insights_generator_input_tokens=0,
            insights_generator_output_tokens=0,

            # Insights Judge
            insights_judge_reasoning=[],
            insights_judge_input_tokens=0,
            insights_judge_output_tokens=0

    )

    initial_state = default_initial_state.copy()
    llm_config_prefixes = {
        "llm_sql_writer_"        : "llm_config_sql_writer",
        "llm_sql_validator_"     : "llm_config_sql_validator",
        "llm_sql_judge_"         : "llm_config_sql_judge",
        "llm_bi_expert_"         : "llm_config_bi_expert",
        "llm_bi_judge_"          : "llm_config_bi_judge",
        "llm_python_writer_"     : "llm_config_python_writer",
        "llm_python_validator_"  : "llm_config_python_validator",
        "llm_python_judge_"      : "llm_config_python_judge",
        "llm_insights_generator_": "llm_config_insights_generator",
        "llm_insights_judge_"    : "llm_config_insights_judge",
    }

    print("Processing kwargs for initial state and LLM configurations...")
    # Separate kwargs for direct state override vs. LLM param override
    llm_param_kwargs = {}
    state_override_kwargs = {}

    for key, value in kwargs.items():
        is_llm_param = False
        for prefix, config_key in llm_config_prefixes.items():
            if key.startswith(prefix):
                if config_key not in llm_param_kwargs:
                    llm_param_kwargs[config_key] = {}
                param_name = key[len(prefix):]  # e.g., "model_name", "temperature"
                llm_param_kwargs[config_key][param_name] = value
                print(f"  - LLM Config: Setting {param_name} for {config_key} to {value}")
                is_llm_param = True
                break
        if not is_llm_param:
            state_override_kwargs[key] = value

    # Apply direct state overrides
    for key, value in state_override_kwargs.items():
        if key in initial_state:
            initial_state[key] = value
            print(f"  - State Override: Setting '{key}' to {value}")
        else:
            print(f"  - Warning: Ignored unknown direct state kwarg '{key}'")

    # Update LLM configurations in the initial_state
    for config_key, params_to_update in llm_param_kwargs.items():
        if config_key in initial_state and isinstance(initial_state[config_key], dict):
            initial_state[config_key].update(params_to_update)
        else:
            print(f"  - Warning: Could not find LLM config key '{config_key}' in initial state to update.")

    if 'max_num_retries_debug_python' not in initial_state:  # type: ignore
        initial_state['max_num_retries_debug_python'] = initial_state['max_num_retries_debug']  # type: ignore

    print("\nFinal Initial State (including LLM configs):")
    import pprint;
    pprint.pprint(initial_state);
    print("-" * 20)

    final_state = app.invoke(initial_state)
    return final_state
